"""
Multi-format document ingest.

Dispatches by file extension to a format-specific text extractor. All
optional dependencies are imported lazily so the platform boots even if
only some formats are needed. Missing libraries surface as a clear
'IngestUnsupported' error with the pip install command.

Supported:
    .pdf                         -> pypdf
    .docx                        -> python-docx
    .pptx                        -> python-pptx
    .xlsx, .xlsm                 -> openpyxl
    .csv, .tsv                   -> stdlib csv
    .txt, .md, .markdown, .log   -> stdlib (utf-8 with fallback)
    .html, .htm                  -> stdlib HTMLParser
    .rtf                         -> striprtf
    .json                        -> stdlib json (pretty-printed back to text)

Each extractor returns a single string. The caller is responsible for
chunking and downstream LLM extraction.
"""

from __future__ import annotations

import csv
import hashlib
import io
import json
from html.parser import HTMLParser
from pathlib import Path
from typing import Callable

SUPPORTED_EXTS: dict[str, str] = {
    ".pdf":      "pdf",
    ".docx":     "docx",
    ".pptx":     "pptx",
    ".xlsx":     "xlsx",
    ".xlsm":     "xlsx",
    ".csv":      "csv",
    ".tsv":      "tsv",
    ".txt":      "text",
    ".md":       "text",
    ".markdown": "text",
    ".log":      "text",
    ".html":     "html",
    ".htm":      "html",
    ".rtf":      "rtf",
    ".json":     "json",
}


class IngestUnsupported(RuntimeError):
    """Raised when a file's extension isn't handled or a needed lib is missing."""


def content_hash(path: Path) -> str:
    h = hashlib.sha256()
    with path.open("rb") as fh:
        for chunk in iter(lambda: fh.read(65536), b""):
            h.update(chunk)
    return h.hexdigest()


# ----------------------------------------------------------------------
# Per-format extractors
# ----------------------------------------------------------------------

def _extract_pdf(path: Path) -> str:
    try:
        from pypdf import PdfReader
    except ImportError as e:
        raise IngestUnsupported("pypdf required: pip install pypdf") from e
    reader = PdfReader(str(path))
    return "\n\n".join((p.extract_text() or "") for p in reader.pages)


def _extract_docx(path: Path) -> str:
    try:
        import docx  # python-docx
    except ImportError as e:
        raise IngestUnsupported("python-docx required: pip install python-docx") from e
    d = docx.Document(str(path))
    parts: list[str] = []
    for para in d.paragraphs:
        if para.text.strip():
            parts.append(para.text)
    for table in d.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells]
            if any(cells):
                parts.append("\t".join(cells))
    return "\n".join(parts)


def _extract_pptx(path: Path) -> str:
    try:
        from pptx import Presentation
    except ImportError as e:
        raise IngestUnsupported("python-pptx required: pip install python-pptx") from e
    prs = Presentation(str(path))
    out: list[str] = []
    for i, slide in enumerate(prs.slides, start=1):
        out.append(f"--- Slide {i} ---")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in para.runs)
                    if text.strip():
                        out.append(text)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells]
                    out.append("\t".join(cells))
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text
            if notes.strip():
                out.append(f"[notes] {notes}")
    return "\n".join(out)


def _extract_xlsx(path: Path) -> str:
    try:
        from openpyxl import load_workbook
    except ImportError as e:
        raise IngestUnsupported("openpyxl required: pip install openpyxl") from e
    wb = load_workbook(filename=str(path), read_only=True, data_only=True)
    out: list[str] = []
    for ws in wb.worksheets:
        out.append(f"--- Sheet: {ws.title} ---")
        for row in ws.iter_rows(values_only=True):
            cells = ["" if v is None else str(v) for v in row]
            if any(c.strip() for c in cells):
                out.append("\t".join(cells))
    return "\n".join(out)


def _extract_csv(path: Path, *, delim: str = ",") -> str:
    out: list[str] = []
    with path.open("r", encoding="utf-8", errors="replace", newline="") as fh:
        for row in csv.reader(fh, delimiter=delim):
            out.append("\t".join(row))
    return "\n".join(out)


def _extract_tsv(path: Path) -> str:
    return _extract_csv(path, delim="\t")


def _extract_text(path: Path) -> str:
    try:
        return path.read_text(encoding="utf-8")
    except UnicodeDecodeError:
        return path.read_text(encoding="latin-1", errors="replace")


class _HtmlText(HTMLParser):
    def __init__(self) -> None:
        super().__init__()
        self._buf: list[str] = []
        self._skip = 0

    def handle_starttag(self, tag, attrs):
        if tag in ("script", "style"):
            self._skip += 1

    def handle_endtag(self, tag):
        if tag in ("script", "style") and self._skip:
            self._skip -= 1

    def handle_data(self, data):
        if self._skip:
            return
        s = data.strip()
        if s:
            self._buf.append(s)

    def text(self) -> str:
        return "\n".join(self._buf)


def _extract_html(path: Path) -> str:
    p = _HtmlText()
    p.feed(_extract_text(path))
    return p.text()


def _extract_rtf(path: Path) -> str:
    try:
        from striprtf.striprtf import rtf_to_text  # type: ignore
    except ImportError as e:
        raise IngestUnsupported("striprtf required: pip install striprtf") from e
    return rtf_to_text(_extract_text(path))


def _extract_json(path: Path) -> str:
    try:
        data = json.loads(_extract_text(path))
    except Exception:
        return _extract_text(path)
    return json.dumps(data, indent=2, ensure_ascii=False)


_DISPATCH: dict[str, Callable[[Path], str]] = {
    "pdf": _extract_pdf,
    "docx": _extract_docx,
    "pptx": _extract_pptx,
    "xlsx": _extract_xlsx,
    "csv": _extract_csv,
    "tsv": _extract_tsv,
    "text": _extract_text,
    "html": _extract_html,
    "rtf": _extract_rtf,
    "json": _extract_json,
}


def is_supported(path: Path) -> bool:
    return path.suffix.lower() in SUPPORTED_EXTS


def extract_text(path: Path) -> str:
    """Return plain text for any supported document format."""
    ext = path.suffix.lower()
    fmt = SUPPORTED_EXTS.get(ext)
    if fmt is None:
        raise IngestUnsupported(
            f"Unsupported extension '{ext}'. Supported: "
            + ", ".join(sorted(SUPPORTED_EXTS))
        )
    return _DISPATCH[fmt](path)


def iter_documents(corpus_dir: Path):
    """Yield (path, content_hash, text) for every supported file in corpus_dir."""
    for p in sorted(corpus_dir.iterdir()):
        if p.is_dir() or p.name.startswith(".") or p.name == "README.md":
            continue
        if not is_supported(p):
            continue
        yield p, content_hash(p), extract_text(p)


if __name__ == "__main__":
    import sys
    path = Path(sys.argv[1])
    print(extract_text(path)[:2000])
